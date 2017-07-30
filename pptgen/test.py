import os
from collections import OrderedDict

import wikipedia
from pptx import Presentation

#0 Title (presentation title slide)
#1 Title and Content
#2 Section Header (sometimes called Segue)
#3 Two Content (side by side bullet textboxes)
#4 Comparison (same but additional title for each side by side content box)
#5 Title Only
#6 Blank
#7 Content with Caption
#8 Picture with Caption

EXCLUDE_H2 = ['Contents', 'References', 'External links', 'See also', 'Bibliography', 'Further reading']

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

def printInfo(query):
    print('Looking for ' + query + ' in Wikipedia...')
    whatever = wikipedia.page(query)

    print('Making a list of sections for ' + query + '...')
    allsections = [sec for sec in whatever.sections if sec not in EXCLUDE_H2]

    # print(akpage.summary)

    topass = OrderedDict()
    topass['Introduction'] = whatever.summary

    print('Building a dictionary for all the sections and corresponding data...')
    for sec in allsections:
        if(whatever.section(sec) != ''):
            print('Found ' + sec)
            topass[sec] = whatever.section(sec)

    return topass

def test(query):
    prs = Presentation()

    topass = printInfo(query)
    # pprint.pprint(topass)
    alltitles = [title for title in topass.keys()]

    for i in range(len(alltitles)):
        if topass[alltitles[i]] == None:
            continue
        #Select 'Bulleted Layout(1) layout for the slide'
        bullet_slide_layout = prs.slide_layouts[1]
        #Add a slide to the current presentation
        slide = prs.slides.add_slide(bullet_slide_layout)

        #Get all the shapes from the slide
        shapes = slide.shapes

        #Get title's shape and body's shape
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        #Assign the heading to the title's shape
        title_shape.text = alltitles[i]

        stufftoadd = list(topass[alltitles[i]].split('. '))

        #Get the text frame for body's shape(content)
        tf = body_shape.text_frame
        #tf.text = topass[alltitles[i]]

        i = 0

        for line in stufftoadd:
            if i > 1:
                break
            p = tf.add_paragraph()
            p.level = 0
            p.text = line
            i += 1

    print('Saving the ppt file...')
    #Save file/ Overwrite file
    prs.save(BASE_DIR + '/pptgen/static/pptgen/PPTS/' + query + '.pptx')

if __name__ == '__main__':
    test(input('Enter query>>>'))