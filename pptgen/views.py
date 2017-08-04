from django.shortcuts import render

import os
from collections import OrderedDict

import wikipedia
from pptx import Presentation

def main(request):
    if request.method == 'GET':
        return render(request, 'pptgen/main.html')
    else:
        EXCLUDE_H2 = ['Contents', 'References', 'External links', 'See also', 'Bibliography', 'Further reading']

        BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

        prs = Presentation()

        query = request.POST['ppt_topic']

        print('Looking for ' + query + ' in Wikipedia...')
        whatever = wikipedia.page(query)

        print('Making a list of sections for ' + query + '...')
        allsections = [sec for sec in whatever.sections if sec not in EXCLUDE_H2]

        # print(akpage.summary)

        topass = OrderedDict()
        topass['Introduction'] = whatever.summary

        print('Building a dictionary for all the sections and corresponding data...')
        for sec in allsections:
            if(whatever.section(sec) != ''):
                print('Found ' + sec)
                topass[sec] = whatever.section(sec)

        # pprint.pprint(topass)
        alltitles = [title for title in topass.keys()]

        for i in range(len(alltitles)):
            if topass[alltitles[i]] == None:
                continue
            #Select 'Bulleted Layout(1) layout for the slide'
            bullet_slide_layout = prs.slide_layouts[1]
            #Add a slide to the current presentation
            slide = prs.slides.add_slide(bullet_slide_layout)

            #Get all the shapes from the slide
            shapes = slide.shapes

            #Get title's shape and body's shape
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            #Assign the heading to the title's shape
            title_shape.text = alltitles[i]

            stufftoadd = list(topass[alltitles[i]].split('. '))

            #Get the text frame for body's shape(content)
            tf = body_shape.text_frame
            #tf.text = topass[alltitles[i]]

            i = 0

            for line in stufftoadd:
                if i > 1:
                    break
                p = tf.add_paragraph()
                p.level = 0
                p.text = line
                i += 1

        print('Saving the ppt file...')
        #Save file/ Overwrite file
        prs.save(BASE_DIR + '/pptgen/static/pptgen/PPTS/' + query + '.pptx')

        PPT_DIR = BASE_DIR + '/pptgen/PPTS'
        return render(request, 'pptgen/generated.html', {'ppt_topic': query, 'PPT_DIR': PPT_DIR})